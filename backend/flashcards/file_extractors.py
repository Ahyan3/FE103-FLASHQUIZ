"""
File text extraction utilities for PDF, DOCX, and PPTX files
"""
import io
import PyPDF2
from docx import Document
from pptx import Presentation


def extract_text_from_pdf(file_obj):
    """
    Extract text from PDF file
    
    Args:
        file_obj: File object or BytesIO object
        
    Returns:
        str: Extracted text content
    """
    try:
        pdf_reader = PyPDF2.PdfReader(file_obj)
        text_content = []
        
        for page in pdf_reader.pages:
            text = page.extract_text()
            if text.strip():
                text_content.append(text)
        
        return "\n\n".join(text_content)
    except Exception as e:
        raise ValueError(f"Error extracting text from PDF: {str(e)}")


def extract_text_from_docx(file_obj):
    """
    Extract text from DOCX file
    
    Args:
        file_obj: File object or BytesIO object
        
    Returns:
        str: Extracted text content
    """
    try:
        doc = Document(file_obj)
        text_content = []
        
        # Extract text from paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text)
                if row_text:
                    text_content.append(" | ".join(row_text))
        
        return "\n\n".join(text_content)
    except Exception as e:
        raise ValueError(f"Error extracting text from DOCX: {str(e)}")


def extract_text_from_pptx(file_obj):
    """
    Extract text from PPTX file
    
    Args:
        file_obj: File object or BytesIO object
        
    Returns:
        str: Extracted text content
    """
    try:
        prs = Presentation(file_obj)
        text_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if len(slide_text) > 1:  # Has content beyond title
                text_content.append("\n".join(slide_text))
        
        return "\n\n".join(text_content)
    except Exception as e:
        raise ValueError(f"Error extracting text from PPTX: {str(e)}")


def extract_text_from_file(file_obj, file_type):
    """
    Extract text from supported file types
    
    Args:
        file_obj: File object or BytesIO object
        file_type: str, one of 'pdf', 'docx', 'pptx'
        
    Returns:
        str: Extracted text content
        
    Raises:
        ValueError: If file type is not supported or extraction fails
    """
    file_type = file_type.lower()
    
    extractors = {
        'pdf': extract_text_from_pdf,
        'docx': extract_text_from_docx,
        'pptx': extract_text_from_pptx,
    }
    
    if file_type not in extractors:
        raise ValueError(f"Unsupported file type: {file_type}. Supported types: {', '.join(extractors.keys())}")
    
    return extractors[file_type](file_obj)


def get_file_type_from_filename(filename):
    """
    Get file type from filename extension
    
    Args:
        filename: str, the filename
        
    Returns:
        str: File type (pdf, docx, pptx) or None if unsupported
    """
    filename = filename.lower()
    
    if filename.endswith('.pdf'):
        return 'pdf'
    elif filename.endswith('.docx'):
        return 'docx'
    elif filename.endswith('.pptx'):
        return 'pptx'
    
    return None